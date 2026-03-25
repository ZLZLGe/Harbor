from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/root")
ASSETS_DIR = ROOT / "proposal-assets"
TEMPLATE_PATH = ROOT / "GreenGrid-Template-Workbook.pptx"
BRIEF_PATH = ROOT / "proposal_brief.yaml"
TMP_SLOT = ROOT / "_slot-placeholder.png"

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

PRIMARY = "1F6B5C"
PRIMARY_DARK = "153D35"
ACCENT = "8CC63F"
SAND = "F4F1E8"
INK = "24312E"
MUTED = "6C7A76"
MINT = "E3EFEA"
WARM = "D97C3B"
WHITE = "FFFFFF"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def make_scene_image(path: Path, *, title: str, subtitle: str, band: str, accent: str) -> None:
    image = Image.new("RGB", (1600, 900), hex_to_rgb("F8F6F0"))
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 480, 900), fill=hex_to_rgb(band))
    draw.rectangle((980, 0, 1600, 900), fill=hex_to_rgb(accent))
    draw.rounded_rectangle((120, 120, 1480, 780), radius=42, fill=hex_to_rgb("FFFFFF"))
    draw.rounded_rectangle((1050, 150, 1410, 390), radius=28, fill=hex_to_rgb("EAF3D8"))
    draw.rounded_rectangle((1050, 430, 1410, 670), radius=28, fill=hex_to_rgb("E5F2ED"))
    draw.line((240, 640, 920, 640), fill=hex_to_rgb("BED6CF"), width=10)
    draw.line((240, 690, 780, 690), fill=hex_to_rgb("BED6CF"), width=10)
    draw.text((180, 190), title, fill=hex_to_rgb(PRIMARY_DARK), font=font(62, bold=True))
    draw.text((180, 305), subtitle, fill=hex_to_rgb("516360"), font=font(30))
    draw.text((180, 710), "GreenGrid proposal asset", fill=hex_to_rgb("7A8A84"), font=font(24))
    image.save(path)


def make_placeholder_image(path: Path) -> None:
    image = Image.new("RGB", (1400, 900), hex_to_rgb("DCE9E4"))
    draw = ImageDraw.Draw(image)
    draw.rectangle((60, 60, 1340, 840), outline=hex_to_rgb(PRIMARY), width=14)
    draw.line((60, 60, 1340, 840), fill=hex_to_rgb(PRIMARY), width=10)
    draw.line((1340, 60, 60, 840), fill=hex_to_rgb(PRIMARY), width=10)
    draw.text((280, 390), "[replace image]", fill=hex_to_rgb(PRIMARY_DARK), font=font(54, bold=True))
    image.save(path)


def style_run(run, *, size: int, color: str, bold: bool = False, italic: bool = False) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic


def textbox(slide, left: float, top: float, width: float, height: float, text: str, *, size: int, color: str, bold: bool = False, italic: bool = False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    style_run(run, size=size, color=color, bold=bold, italic=italic)
    return shape


def bullet_box(slide, left: float, top: float, width: float, height: float, items: list[str]) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(8)
        for run in paragraph.runs:
            style_run(run, size=16, color=INK)


def image_slot(slide, slot_name: str, *, left: float, top: float, width: float, height: float):
    picture = slide.shapes.add_picture(str(TMP_SLOT), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    picture._element.nvPicPr.cNvPr.set("name", slot_name)
    return picture


def panel(slide, *, left: float, top: float, width: float, height: float, fill: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()


def base_slide(prs: Presentation, background: str = SAND):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(background)
    return slide


def add_header_band(slide, title: str) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.85))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(PRIMARY)
    band.line.fill.background()
    textbox(slide, 0.72, 0.18, 6.8, 0.38, title, size=26, color=WHITE, bold=True)


def build_template() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = base_slide(prs)
    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(4.0), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = rgb(PRIMARY)
    side.line.fill.background()
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.9), Inches(2.3), Inches(0.42))
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(ACCENT)
    chip.line.fill.background()
    textbox(slide, 0.88, 0.98, 1.95, 0.2, "{{COVER_KICKER}}", size=13, color=WHITE, bold=True)
    textbox(slide, 0.72, 1.55, 4.9, 1.5, "{{COVER_TITLE}}", size=28, color=WHITE, bold=True)
    textbox(slide, 0.72, 3.35, 4.6, 0.7, "{{COVER_SUBTITLE}}", size=18, color=WHITE)
    textbox(slide, 0.72, 4.2, 4.8, 0.9, "{{COVER_TAGLINE}}", size=15, color="E8F0EC", italic=True)
    image_slot(slide, "slot-cover-image", left=6.2, top=0.72, width=6.4, height=5.75)
    textbox(slide, 6.4, 6.7, 5.8, 0.26, "[replace with client-ready cover image]", size=11, color=MUTED, italic=True)

    slide = base_slide(prs, background=MINT)
    add_header_band(slide, "Template divider - remove unless needed")
    textbox(slide, 0.9, 1.55, 9.6, 1.0, "A spare section divider is included here as a template option.", size=28, color=PRIMARY_DARK, bold=True)
    textbox(slide, 0.9, 2.8, 8.8, 0.5, "[Delete this divider slide for the final client deck]", size=16, color=WARM, italic=True)

    slide = base_slide(prs)
    add_header_band(slide, "{{OPPORTUNITY_TITLE}}")
    textbox(slide, 0.85, 1.25, 5.4, 1.0, "{{OPPORTUNITY_INTRO}}", size=20, color=INK)
    bullet_box(slide, 0.95, 2.35, 5.0, 1.6, ["{{OPPORTUNITY_POINT_1}}", "{{OPPORTUNITY_POINT_2}}"])
    panel(slide, left=0.92, top=4.55, width=2.2, height=1.55, fill=PRIMARY)
    textbox(slide, 1.18, 4.88, 1.62, 0.32, "{{OPPORTUNITY_STAT_LABEL}}", size=12, color="D9ECE7", bold=True)
    textbox(slide, 1.18, 5.28, 1.62, 0.42, "{{OPPORTUNITY_STAT_VALUE}}", size=24, color=WHITE, bold=True)
    image_slot(slide, "slot-opportunity-image", left=6.55, top=1.28, width=5.9, height=4.9)
    textbox(slide, 6.65, 6.32, 4.9, 0.25, "[replace with opportunity image]", size=11, color=MUTED, italic=True)

    slide = base_slide(prs, background=MINT)
    add_header_band(slide, "{{SOLUTION_TITLE}}")
    textbox(slide, 0.9, 1.22, 10.8, 0.85, "{{SOLUTION_INTRO}}", size=19, color=INK)
    card_lefts = [0.92, 4.47, 8.02]
    for index, left in enumerate(card_lefts, start=1):
        panel(slide, left=left, top=2.3, width=3.0, height=2.35, fill=WHITE)
        textbox(slide, left + 0.22, 2.66, 2.56, 1.4, "{{SOLUTION_PILLAR_%d}}" % index, size=18, color=PRIMARY_DARK, bold=True)
    textbox(slide, 0.94, 5.35, 10.5, 0.65, "{{SOLUTION_FOOTER}}", size=14, color=MUTED, italic=True)

    slide = base_slide(prs)
    add_header_band(slide, "{{PILOT_TITLE}}")
    phase_lefts = [0.85, 4.48, 8.11]
    for index, left in enumerate(phase_lefts, start=1):
        panel(slide, left=left, top=1.55, width=3.0, height=4.15, fill=WHITE)
        textbox(slide, left + 0.22, 1.92, 2.4, 0.3, "{{PILOT_WINDOW_%d}}" % index, size=13, color=PRIMARY, bold=True)
        textbox(slide, left + 0.22, 2.45, 2.45, 0.85, "{{PILOT_HEADING_%d}}" % index, size=20, color=INK, bold=True)
        textbox(slide, left + 0.22, 3.55, 2.45, 1.25, "{{PILOT_DETAIL_%d}}" % index, size=15, color=MUTED)

    slide = base_slide(prs, background=MINT)
    add_header_band(slide, "Spare metrics board")
    textbox(slide, 0.95, 1.45, 8.8, 0.8, "This board is available in the template, but it is not required for this client proposal.", size=24, color=PRIMARY_DARK, bold=True)
    textbox(slide, 0.95, 2.65, 8.3, 0.4, "[Unused template option - remove in final output]", size=16, color=WARM, italic=True)

    slide = base_slide(prs)
    add_header_band(slide, "{{PROOF_TITLE}}")
    panel(slide, left=0.92, top=1.4, width=5.3, height=3.1, fill=WHITE)
    textbox(slide, 1.2, 1.8, 4.6, 1.55, "{{PROOF_QUOTE}}", size=22, color=PRIMARY_DARK, italic=True)
    textbox(slide, 1.2, 3.55, 4.6, 0.32, "{{PROOF_ATTRIBUTION}}", size=13, color=MUTED, bold=True)
    panel(slide, left=0.92, top=5.0, width=2.4, height=1.2, fill=PRIMARY)
    textbox(slide, 1.18, 5.18, 1.9, 0.42, "{{PROOF_VALUE_1}}", size=24, color=WHITE, bold=True)
    textbox(slide, 1.18, 5.62, 1.9, 0.26, "{{PROOF_LABEL_1}}", size=11, color="D9ECE7")
    panel(slide, left=3.58, top=5.0, width=2.4, height=1.2, fill=ACCENT)
    textbox(slide, 3.85, 5.18, 1.8, 0.42, "{{PROOF_VALUE_2}}", size=24, color=PRIMARY_DARK, bold=True)
    textbox(slide, 3.85, 5.62, 1.8, 0.26, "{{PROOF_LABEL_2}}", size=11, color=PRIMARY_DARK)
    image_slot(slide, "slot-proof-image", left=6.7, top=1.42, width=5.75, height=4.82)
    textbox(slide, 6.85, 6.36, 4.4, 0.25, "[replace with proof image]", size=11, color=MUTED, italic=True)

    slide = base_slide(prs, background=MINT)
    add_header_band(slide, "{{NEXT_TITLE}}")
    bullet_box(slide, 1.0, 1.55, 8.9, 2.55, ["{{NEXT_STEP_1}}", "{{NEXT_STEP_2}}", "{{NEXT_STEP_3}}"])
    panel(slide, left=0.95, top=4.65, width=11.2, height=1.05, fill=WHITE)
    textbox(slide, 1.2, 4.98, 10.7, 0.32, "{{NEXT_FOOTER}}", size=14, color=PRIMARY_DARK)

    prs.save(TEMPLATE_PATH)


def patch_theme() -> None:
    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)
        with zipfile.ZipFile(TEMPLATE_PATH) as archive:
            archive.extractall(tmpdir_path)

        theme_path = tmpdir_path / "ppt" / "theme" / "theme1.xml"
        root = ET.parse(theme_path).getroot()
        clr_scheme = root.find(".//a:clrScheme", NS)
        colors = {
            "accent1": PRIMARY,
            "accent2": ACCENT,
            "accent3": "BED6CF",
            "accent4": WARM,
            "accent5": "4C7F74",
            "accent6": "A8B89F",
            "hlink": PRIMARY_DARK,
            "folHlink": "5A7C72",
        }
        if clr_scheme is not None:
            for key, value in colors.items():
                node = clr_scheme.find(f"a:{key}/a:srgbClr", NS)
                if node is not None:
                    node.set("val", value)

        major_latin = root.find(".//a:fontScheme/a:majorFont/a:latin", NS)
        minor_latin = root.find(".//a:fontScheme/a:minorFont/a:latin", NS)
        if major_latin is not None:
            major_latin.set("typeface", "Trebuchet MS")
        if minor_latin is not None:
            minor_latin.set("typeface", "Verdana")

        ET.register_namespace("a", NS["a"])
        ET.ElementTree(root).write(theme_path, encoding="utf-8", xml_declaration=True)

        rebuilt = tmpdir_path / "rebuilt.pptx"
        with zipfile.ZipFile(rebuilt, "w", zipfile.ZIP_DEFLATED) as archive:
            for path in sorted(tmpdir_path.rglob("*")):
                if path == rebuilt or path.is_dir():
                    continue
                archive.write(path, path.relative_to(tmpdir_path).as_posix())
        shutil.move(rebuilt, TEMPLATE_PATH)


def main() -> None:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    make_scene_image(
        ASSETS_DIR / "site-aerial.png",
        title="Phoenix cross-dock",
        subtitle="High-load yard and canopy opportunity",
        band=PRIMARY,
        accent="DCE8B6",
    )
    make_scene_image(
        ASSETS_DIR / "operations-dashboard.png",
        title="Dispatch dashboard",
        subtitle="Charge state, tariff window, and savings signals",
        band=PRIMARY_DARK,
        accent="E6F2EE",
    )
    make_scene_image(
        ASSETS_DIR / "team-workshop.png",
        title="Operator workshop",
        subtitle="Pilot planning with site operations and finance",
        band=WARM,
        accent="F4E6D9",
    )
    make_placeholder_image(TMP_SLOT)
    build_template()
    patch_theme()
    TMP_SLOT.unlink(missing_ok=True)
    if not BRIEF_PATH.exists():
        raise FileNotFoundError(BRIEF_PATH)


if __name__ == "__main__":
    main()
