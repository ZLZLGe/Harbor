from __future__ import annotations

import json
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

BRIEF_PATH = Path("/root/brief.md")
METRICS_PATH = Path("/root/metrics.json")
PALETTE_PATH = Path("/root/brand/palette.json")
OUTPUT_PPTX = Path("/root/Support-Onboarding-playbook.pptx")

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

THEME_COLOR_TO_SCHEME = {
    "ACCENT_1": "accent1",
    "ACCENT_2": "accent2",
    "ACCENT_3": "accent3",
    "ACCENT_4": "accent4",
    "ACCENT_5": "accent5",
    "ACCENT_6": "accent6",
    "BACKGROUND_1": "bg1",
    "BACKGROUND_2": "bg2",
    "DARK_1": "dk1",
    "DARK_2": "dk2",
    "FOLLOWED_HYPERLINK": "folHlink",
    "HYPERLINK": "hlink",
    "LIGHT_1": "lt1",
    "LIGHT_2": "lt2",
    "TEXT_1": "tx1",
    "TEXT_2": "tx2",
}


def parse_brief(path: Path) -> dict[str, object]:
    text = path.read_text(encoding="utf-8")

    def section(name: str) -> str:
        pattern = rf"## {re.escape(name)}\n(.*?)(?=\n## |\Z)"
        match = re.search(pattern, text, re.S)
        assert match is not None, f"missing section {name}"
        return match.group(1).strip()

    cover = section("Cover")
    workflow = section("Workflow")
    closing = section("Closing")

    return {
        "cover_title": re.search(r"Title:\s*(.+)", cover).group(1).strip(),
        "cover_subtitle": re.search(r"Subtitle:\s*(.+)", cover).group(1).strip(),
        "cover_tagline": re.search(r"Tagline:\s*(.+)", cover).group(1).strip(),
        "agenda": [match.strip() for match in re.findall(r"^\d+\.\s+(.+)$", section("Agenda"), re.M)],
        "workflow_title": re.search(r"Title:\s*(.+)", workflow).group(1).strip(),
        "workflow_intro": re.search(r"Intro:\s*(.+)", workflow).group(1).strip(),
        "workflow_steps": [tuple(part.strip() for part in line[2:].split("|")) for line in workflow.splitlines() if line.startswith("- ")],
        "closing_title": re.search(r"Title:\s*(.+)", closing).group(1).strip(),
        "closing_bullets": [match.strip() for match in re.findall(r"^- (.+)$", closing, re.M)],
        "closing_footer": re.search(r"Footer:\s*(.+)", closing).group(1).strip(),
    }


def load_presentation() -> Presentation:
    return Presentation(str(OUTPUT_PPTX))


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def iter_paragraphs(slide):
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            yield from shape.text_frame.paragraphs
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield from cell.text_frame.paragraphs


def normalize(text: str) -> str:
    return " ".join(text.split())


def slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for paragraph in iter_paragraphs(slide):
        text = normalize(paragraph.text)
        if text:
            texts.append(text)
    return texts


def first_table_matrix(slide) -> list[list[str]]:
    for shape in iter_shapes(slide.shapes):
        if not getattr(shape, "has_table", False):
            continue
        rows: list[list[str]] = []
        for row in shape.table.rows:
            cells = [" ".join(filter(None, (normalize(paragraph.text) for paragraph in cell.text_frame.paragraphs))).strip() for cell in row.cells]
            rows.append(cells)
        return rows
    raise AssertionError("expected a table on the slide")


def picture_count(slide) -> int:
    return sum(1 for shape in iter_shapes(slide.shapes) if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def load_theme_colors(path: Path) -> dict[str, str]:
    with zipfile.ZipFile(path) as zipf:
        theme_name = next((name for name in zipf.namelist() if re.fullmatch(r"ppt/theme/theme\d+\.xml", name)), None)
        if theme_name is None:
            return {}
        root = ET.fromstring(zipf.read(theme_name))
        scheme = root.find(".//a:clrScheme", NS)
        if scheme is None:
            return {}

        colors: dict[str, str] = {}
        for child in scheme:
            scheme_name = child.tag.rsplit("}", 1)[-1]
            color_node = next(iter(child), None)
            if color_node is None:
                continue
            color_tag = color_node.tag.rsplit("}", 1)[-1]
            if color_tag == "srgbClr":
                value = color_node.attrib.get("val")
            elif color_tag == "sysClr":
                value = color_node.attrib.get("lastClr")
            else:
                value = None
            if value:
                colors[scheme_name] = value.upper()
        return colors


def resolve_color(color_format, theme_colors: dict[str, str]) -> str | None:
    if color_format is None:
        return None
    if color_format.type == MSO_COLOR_TYPE.RGB:
        return str(color_format.rgb).upper()
    if color_format.type == MSO_COLOR_TYPE.SCHEME:
        theme_color = getattr(color_format, "theme_color", None)
        if theme_color is None:
            return None
        scheme_name = THEME_COLOR_TO_SCHEME.get(theme_color.name)
        if scheme_name is None:
            return None
        return theme_colors.get(scheme_name)
    return None


def shape_fill_color(shape, theme_colors: dict[str, str]) -> str | None:
    fill = getattr(shape, "fill", None)
    if fill is None or fill.type != MSO_FILL.SOLID:
        return None
    return resolve_color(fill.fore_color, theme_colors)


def shape_line_color(shape, theme_colors: dict[str, str]) -> str | None:
    line = getattr(shape, "line", None)
    if line is None:
        return None
    fill = getattr(line, "fill", None)
    if fill is None or fill.type != MSO_FILL.SOLID:
        return None
    return resolve_color(line.color, theme_colors)


def run_color(run, theme_colors: dict[str, str]) -> str | None:
    return resolve_color(run.font.color, theme_colors)


def shape_area(shape) -> int:
    return int(shape.width) * int(shape.height)


def slide_has_canvas_background(slide, canvas: str, slide_area: int, theme_colors: dict[str, str]) -> bool:
    background_fill = slide.background.fill
    if background_fill.type == MSO_FILL.SOLID and resolve_color(background_fill.fore_color, theme_colors) == canvas:
        return True

    largest_canvas_area = 0
    for shape in iter_shapes(slide.shapes):
        if shape_fill_color(shape, theme_colors) == canvas:
            largest_canvas_area = max(largest_canvas_area, shape_area(shape))
    return largest_canvas_area >= slide_area * 0.5


def slide_has_primary_title_treatment(slide, expected_title: str, primary: str, slide_height: int, theme_colors: dict[str, str]) -> bool:
    for paragraph in iter_paragraphs(slide):
        if normalize(paragraph.text) != expected_title:
            continue
        for run in paragraph.runs:
            if normalize(run.text) == expected_title and run_color(run, theme_colors) == primary:
                return True

    for shape in iter_shapes(slide.shapes):
        if int(shape.top) > slide_height * 0.3:
            continue
        if shape_fill_color(shape, theme_colors) == primary or shape_line_color(shape, theme_colors) == primary:
            return True
    return False


def accent_element_count(slide, accent: str, slide_area: int, theme_colors: dict[str, str]) -> int:
    count = 0
    for shape in iter_shapes(slide.shapes):
        if shape_fill_color(shape, theme_colors) == accent and shape_area(shape) <= slide_area * 0.4:
            count += 1
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run_color(run, theme_colors) == accent:
                        count += 1
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run_color(run, theme_colors) == accent:
                                count += 1
    return count


def test_output_exists() -> None:
    assert OUTPUT_PPTX.exists(), "expected output deck at /root/Support-Onboarding-playbook.pptx"
    assert OUTPUT_PPTX.stat().st_size > 0, "output deck is empty"


def test_slide_count_is_exactly_five() -> None:
    presentation = load_presentation()
    assert len(presentation.slides) == 5, "expected exactly 5 slides"


def test_cover_slide_uses_brief_content_and_wordmark() -> None:
    brief = parse_brief(BRIEF_PATH)
    slide = load_presentation().slides[0]
    texts = slide_texts(slide)
    assert brief["cover_title"] in texts, "cover title missing"
    assert brief["cover_subtitle"] in texts, "cover subtitle missing"
    assert brief["cover_tagline"] in texts, "cover tagline missing"
    assert picture_count(slide) >= 1, "cover slide should include the provided wordmark asset"


def test_agenda_slide_contains_ordered_items() -> None:
    brief = parse_brief(BRIEF_PATH)
    slide = load_presentation().slides[1]
    texts = slide_texts(slide)
    assert "Agenda" in texts, "agenda slide title missing"
    positions = []
    for item in brief["agenda"]:
        assert item in texts, f"agenda item missing: {item}"
        positions.append(texts.index(item))
    assert positions == sorted(positions), "agenda items must stay in brief order"


def test_workflow_slide_contains_all_stage_content() -> None:
    brief = parse_brief(BRIEF_PATH)
    slide = load_presentation().slides[2]
    texts = slide_texts(slide)
    assert brief["workflow_title"] in texts, "workflow title missing"
    assert brief["workflow_intro"] in texts, "workflow intro missing"
    for stage, description, owner in brief["workflow_steps"]:
        assert stage in texts, f"workflow stage missing: {stage}"
        assert description in texts, f"workflow description missing for {stage}"
        assert owner in texts, f"workflow owner missing for {stage}"


def test_metrics_slide_has_headline_and_exact_table() -> None:
    metrics = json.loads(METRICS_PATH.read_text(encoding="utf-8"))
    slide = load_presentation().slides[3]
    texts = slide_texts(slide)
    assert metrics["slide_title"] in texts, "metrics slide title missing"
    assert metrics["headline_metric"]["value"] in texts, "headline metric value missing"
    assert metrics["headline_metric"]["label"] in texts, "headline metric label missing"
    assert metrics["headline_metric"]["note"] in texts, "headline metric note missing"

    expected_table = [metrics["table"]["columns"], *metrics["table"]["rows"]]
    assert first_table_matrix(slide) == expected_table, "metrics table does not match metrics.json"


def test_closing_slide_contains_footer_and_wordmark() -> None:
    brief = parse_brief(BRIEF_PATH)
    slide = load_presentation().slides[4]
    texts = slide_texts(slide)
    assert brief["closing_title"] in texts, "closing title missing"
    for bullet in brief["closing_bullets"]:
        assert bullet in texts, f"closing bullet missing: {bullet}"
    assert brief["closing_footer"] in texts, "closing footer missing"
    assert picture_count(slide) >= 1, "closing slide should include the provided wordmark asset"


def test_brand_palette_is_used_consistently() -> None:
    brief = parse_brief(BRIEF_PATH)
    metrics = json.loads(METRICS_PATH.read_text(encoding="utf-8"))
    palette = json.loads(PALETTE_PATH.read_text(encoding="utf-8"))
    presentation = load_presentation()
    theme_colors = load_theme_colors(OUTPUT_PPTX)

    primary = palette["primary"].upper()
    accent = palette["accent"].upper()
    canvas = palette["canvas"].upper()
    slide_area = int(presentation.slide_width) * int(presentation.slide_height)
    slide_height = int(presentation.slide_height)
    expected_titles = [
        str(brief["cover_title"]),
        "Agenda",
        str(brief["workflow_title"]),
        metrics["slide_title"],
        str(brief["closing_title"]),
    ]

    accent_slides = 0
    accent_total = 0
    for index, (slide, title) in enumerate(zip(presentation.slides, expected_titles), start=1):
        assert slide_has_primary_title_treatment(slide, title, primary, slide_height, theme_colors), (
            f"slide {index} should apply the primary brand color to its title treatment"
        )
        assert slide_has_canvas_background(slide, canvas, slide_area, theme_colors), (
            f"slide {index} should use the canvas brand color as its dominant light background tone"
        )

        accent_count = accent_element_count(slide, accent, slide_area, theme_colors)
        accent_total += accent_count
        if accent_count > 0:
            accent_slides += 1

    assert accent_slides >= 3, "accent color should appear on multiple slides for emphasis"
    assert accent_total >= 4, "accent color should be used repeatedly for emphasis across the deck"
