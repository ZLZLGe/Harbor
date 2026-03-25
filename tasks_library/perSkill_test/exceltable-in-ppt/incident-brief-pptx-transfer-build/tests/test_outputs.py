import json
import os
import zipfile

from pptx import Presentation

BRIEF_FILE = "/root/incident-brief.json"
BRAND_FILE = "/root/brand-guide.md"
RESULT_FILE = "/root/results-incident-brief.pptx"

EXPECTED_TITLES = [
    None,
    "Response timeline",
    "48-hour action tracker",
]


with open(BRIEF_FILE, "r", encoding="utf-8") as handle:
    BRIEF = json.load(handle)


def parse_brand_guide(path):
    data = {}
    with open(path, "r", encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.strip()
            if not line.startswith("- ") or ": " not in line:
                continue
            key, value = line[2:].split(": ", 1)
            data[key.strip()] = value.strip().lstrip("#")
    return data


BRAND = parse_brand_guide(BRAND_FILE)


def shape_text(shape):
    if not hasattr(shape, "text_frame"):
        return ""
    parts = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def slide_texts(slide):
    texts = []
    for shape in slide.shapes:
        text = shape_text(shape)
        if text:
            texts.append(text)
    return texts


def find_shape_by_exact_text(slide, expected):
    for shape in slide.shapes:
        if shape_text(shape) == expected:
            return shape
    raise AssertionError(f"Could not find shape with text: {expected}")


def first_run_hex(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                rgb = run.font.color.rgb
                return None if rgb is None else str(rgb)
    return None


def find_first_table(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return shape.table
    raise AssertionError("Slide does not contain a native PowerPoint table")


def test_output_exists_and_is_valid_pptx():
    assert os.path.exists(RESULT_FILE), "Output PPTX was not created"
    assert zipfile.is_zipfile(RESULT_FILE), "Output file is not a valid PPTX archive"


def test_presentation_has_exactly_three_slides():
    prs = Presentation(RESULT_FILE)
    assert len(prs.slides) == 3


def test_cover_contains_required_copy_and_badge():
    prs = Presentation(RESULT_FILE)
    slide = prs.slides[0]
    texts = slide_texts(slide)

    assert BRIEF["report_title"] in texts
    assert f'{BRIEF["severity"]} | {BRIEF["location"]} | {BRIEF["report_date"]}' in texts
    assert BRIEF["executive_summary"] in texts

    badge = find_shape_by_exact_text(slide, BRIEF["severity"])
    fill_rgb = badge.fill.fore_color.rgb
    assert fill_rgb is not None, "Severity badge must use a solid fill color"
    assert str(fill_rgb) == BRAND["Accent badge color"]


def test_slide_titles_use_brand_primary_color():
    prs = Presentation(RESULT_FILE)

    cover_title = find_shape_by_exact_text(prs.slides[0], BRIEF["report_title"])
    assert first_run_hex(cover_title) == BRAND["Primary title color"]

    for index, title in enumerate(EXPECTED_TITLES[1:], start=1):
        shape = find_shape_by_exact_text(prs.slides[index], title)
        assert first_run_hex(shape) == BRAND["Primary title color"]


def test_timeline_slide_lists_all_events_in_order():
    prs = Presentation(RESULT_FILE)
    slide = prs.slides[1]
    texts = slide_texts(slide)
    expected_lines = [f'{item["time"]} - {item["event"]}' for item in BRIEF["timeline"]]

    indices = []
    for line in expected_lines:
        assert line in texts, f"Missing timeline item: {line}"
        indices.append(texts.index(line))

    assert indices == sorted(indices), "Timeline items are not in chronological order"


def test_action_table_contains_all_rows_and_headers():
    prs = Presentation(RESULT_FILE)
    slide = prs.slides[2]
    table = find_first_table(slide)

    headers = [table.cell(0, col).text.strip() for col in range(4)]
    assert headers == ["Owner", "Action", "Due", "Status"]
    assert len(table.rows) == len(BRIEF["actions"]) + 1

    actual_rows = []
    for row_idx in range(1, len(table.rows)):
        actual_rows.append(
            {
                "owner": table.cell(row_idx, 0).text.strip(),
                "action": table.cell(row_idx, 1).text.strip(),
                "due": table.cell(row_idx, 2).text.strip(),
                "status": table.cell(row_idx, 3).text.strip(),
            }
        )

    assert actual_rows == BRIEF["actions"]
